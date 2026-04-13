"""File handling — type detection, conversion, and Gemini inline preparation.

Supported natively by Gemini:
  Images: jpeg, png, gif, webp, bmp
  PDFs:   pdf
  Audio:  wav, mp3, m4a, ogg, flac, aiff, webm, aac
  Video:  mp4, mov, avi, mkv
  Text:   txt, md, csv, html, json, xml, py, js, ts, css, sh, yaml, go, rs, java...

Converted before sending:
  DOCX / DOC   → text/plain  (python-docx)
  XLSX / XLS   → text/csv   (openpyxl)
  PPTX / PPT   → text/plain  (python-pptx)
"""

import base64
import json
import mimetypes
import os
import time
import uuid


def _data_dir():
    return os.environ.get("JAIKA_DATA_DIR", "./data")


# ── Extension → MIME map ──────────────────────────────────────────────────────

EXT_MIME: dict[str, str] = {
    # Images
    "jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
    "gif": "image/gif", "webp": "image/webp", "bmp": "image/bmp",
    "ico": "image/x-icon", "tiff": "image/tiff", "tif": "image/tiff",
    # PDF
    "pdf": "application/pdf",
    # Audio
    "wav": "audio/wav", "mp3": "audio/mpeg", "m4a": "audio/mp4",
    "ogg": "audio/ogg", "flac": "audio/flac", "aiff": "audio/aiff",
    "aif": "audio/aiff", "webm": "audio/webm", "aac": "audio/aac",
    # Video
    "mp4": "video/mp4", "mov": "video/quicktime", "avi": "video/x-msvideo",
    "mkv": "video/x-matroska", "wmv": "video/x-ms-wmv", "3gp": "video/3gpp",
    # Text / Code (all go as text/plain or specific text/* Gemini supports)
    "txt": "text/plain", "md": "text/plain", "rst": "text/plain",
    "csv": "text/csv",
    "html": "text/html", "htm": "text/html",
    "json": "application/json", "jsonl": "application/json",
    "xml": "text/xml",
    "py": "text/x-python", "pyw": "text/x-python",
    "js": "text/javascript", "mjs": "text/javascript",
    "ts": "text/x-typescript",
    "css": "text/css",
    "sh": "text/x-sh", "bash": "text/x-sh", "zsh": "text/x-sh",
    "yaml": "text/x-yaml", "yml": "text/x-yaml",
    "toml": "text/plain", "ini": "text/plain", "cfg": "text/plain",
    "go": "text/x-go", "rs": "text/x-rust", "java": "text/x-java",
    "cpp": "text/x-c++", "cc": "text/x-c++", "cxx": "text/x-c++",
    "c": "text/x-c", "h": "text/x-c",
    "rb": "text/x-ruby",
    "sql": "text/x-sql",
    "r": "text/x-r",
    "swift": "text/x-swift",
    "kt": "text/x-kotlin",
    "php": "text/x-php",
    "lua": "text/x-lua",
    "scala": "text/x-scala",
    "dart": "text/x-dart",
    "tf": "text/x-terraform", "hcl": "text/x-terraform",
    "dockerfile": "text/plain",
    "makefile": "text/plain",
    # Notebooks
    "ipynb": "application/x-ipynb+json",
    # Office — need conversion
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "doc":  "application/msword",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "xls":  "application/vnd.ms-excel",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "ppt":  "application/vnd.ms-powerpoint",
}

# MIME types Gemini natively handles as inline_data
GEMINI_NATIVE: set[str] = {
    "image/jpeg", "image/png", "image/gif", "image/webp", "image/bmp",
    "image/x-icon", "image/tiff",
    "application/pdf",
    "audio/wav", "audio/mpeg", "audio/mp4", "audio/ogg", "audio/flac",
    "audio/aiff", "audio/webm", "audio/aac",
    "video/mp4", "video/quicktime", "video/x-msvideo", "video/x-matroska",
    "video/x-ms-wmv", "video/3gpp",
    "text/plain", "text/html", "text/css", "text/javascript",
    "text/x-typescript", "text/x-python", "text/csv", "application/json",
    "text/xml", "text/x-yaml", "text/x-go", "text/x-rust", "text/x-java",
    "text/x-c++", "text/x-c", "text/x-ruby", "text/x-sql", "text/x-r",
    "text/x-swift", "text/x-kotlin", "text/x-php", "text/x-lua",
    "text/x-scala", "text/x-dart", "text/x-sh", "text/x-terraform",
}

NOTEBOOK_MIMES: set[str] = {"application/x-ipynb+json"}

OFFICE_MIMES: set[str] = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}


def get_mime(filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    return EXT_MIME.get(ext) or mimetypes.guess_type(filename)[0] or "application/octet-stream"


def is_supported(filename: str) -> bool:
    """All files are accepted — unknown types get a text-extraction attempt."""
    return True  # we try to handle everything; worst case we send as-is


# ── Upload & storage ──────────────────────────────────────────────────────────

def save_upload(user_id: str, file_obj, filename: str) -> dict:
    """Save an uploaded file, converting office formats to text.

    Returns metadata dict:
      {file_id, name, mime_type, original_type, converted, size, path, created}
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "bin"
    mime = get_mime(filename)
    file_id = uuid.uuid4().hex[:12]

    uploads_dir = os.path.join(_data_dir(), "users", user_id, "uploads")
    os.makedirs(uploads_dir, exist_ok=True)

    raw_path = os.path.join(uploads_dir, f"{file_id}.{ext}")
    file_obj.save(raw_path)
    raw_size = os.path.getsize(raw_path)

    converted = False
    final_path = raw_path
    final_mime = mime

    converted_text = None
    if mime in NOTEBOOK_MIMES:
        converted_text = _ipynb_to_text(raw_path)
        conv_mime = "text/plain"
    elif mime in OFFICE_MIMES:
        converted_text, conv_mime = _convert_office(raw_path, mime)

    if converted_text is not None:
        conv_ext = "csv" if conv_mime == "text/csv" else "txt"
        final_path = os.path.join(uploads_dir, f"{file_id}.{conv_ext}")
        with open(final_path, "w", encoding="utf-8") as f:
            f.write(converted_text)
        try:
            os.remove(raw_path)
        except OSError:
            pass
        final_mime = conv_mime
        converted = True
    elif mime not in GEMINI_NATIVE:
        # Unknown binary — try to read as text, fallback to raw
        text = _try_read_as_text(raw_path)
        if text:
            final_path = os.path.join(uploads_dir, f"{file_id}.txt")
            with open(final_path, "w", encoding="utf-8") as f:
                f.write(text)
            try:
                os.remove(raw_path)
            except OSError:
                pass
            final_mime = "text/plain"
            converted = True

    meta = {
        "file_id": file_id,
        "name": filename,
        "mime_type": final_mime,
        "original_type": mime,
        "converted": converted,
        "size": raw_size,
        "path": final_path,
        "created": time.time(),
    }
    meta_path = os.path.join(uploads_dir, f"{file_id}.meta")
    with open(meta_path, "w") as f:
        json.dump(meta, f)

    return meta


def load_file_inline(file_id: str, user_id: str) -> dict | None:
    """Load file and return as Gemini inline_data dict.

    Returns: {name, mime_type, base64, converted, original_type} or None
    """
    meta = get_file_meta(file_id, user_id)
    if not meta:
        return None

    file_path = meta["path"]
    if not os.path.exists(file_path):
        return None

    with open(file_path, "rb") as f:
        data = f.read()

    return {
        "name": meta["name"],
        "mime_type": meta["mime_type"],
        "base64": base64.b64encode(data).decode("ascii"),
        "converted": meta.get("converted", False),
        "original_type": meta.get("original_type", meta["mime_type"]),
    }


def _safe_id(val):
    """Sanitize an ID to prevent path traversal."""
    return os.path.basename(str(val)).replace("..", "").strip()


def get_file_meta(file_id: str, user_id: str) -> dict | None:
    """Return file metadata without loading content."""
    file_id = _safe_id(file_id)
    uploads_dir = os.path.join(_data_dir(), "users", user_id, "uploads")
    meta_path = os.path.join(uploads_dir, f"{file_id}.meta")
    if not os.path.exists(meta_path):
        return None
    try:
        with open(meta_path) as f:
            return json.load(f)
    except (json.JSONDecodeError, IOError, UnicodeDecodeError):
        return None


def delete_file(file_id: str, user_id: str) -> bool:
    """Delete file and its metadata."""
    meta = get_file_meta(file_id, user_id)
    if not meta:
        return False
    file_path = meta.get("path", "")
    if file_path and os.path.exists(file_path):
        try:
            os.remove(file_path)
        except OSError:
            pass
    uploads_dir = os.path.join(_data_dir(), "users", user_id, "uploads")
    try:
        os.remove(os.path.join(uploads_dir, f"{file_id}.meta"))
    except OSError:
        pass
    return True


def list_user_files(user_id: str) -> list[dict]:
    """Return list of file metadata for a user (no binary data)."""
    uploads_dir = os.path.join(_data_dir(), "users", user_id, "uploads")
    if not os.path.exists(uploads_dir):
        return []
    result = []
    for fname in os.listdir(uploads_dir):
        if not fname.endswith(".meta"):
            continue
        file_id = fname[:-5]
        meta = get_file_meta(file_id, user_id)
        if meta and os.path.exists(meta.get("path", "")):
            result.append({
                "file_id": meta["file_id"],
                "name": meta["name"],
                "mime_type": meta["mime_type"],
                "original_type": meta.get("original_type", meta["mime_type"]),
                "converted": meta.get("converted", False),
                "size": meta.get("size", 0),
                "created": meta.get("created", 0),
            })
    result.sort(key=lambda x: x["created"], reverse=True)
    return result


# ── Office conversion ─────────────────────────────────────────────────────────

def _convert_office(path: str, mime: str) -> tuple[str | None, str]:
    """Convert office file to text. Returns (text, target_mime) or (None, mime)."""
    if "word" in mime or mime == "application/msword":
        return _docx_to_text(path), "text/plain"
    if "spreadsheet" in mime or "excel" in mime:
        return _xlsx_to_csv(path), "text/csv"
    if "presentation" in mime or "powerpoint" in mime:
        return _pptx_to_text(path), "text/plain"
    return None, mime


def _docx_to_text(path: str) -> str | None:
    try:
        from docx import Document
        doc = Document(path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                style = para.style.name if para.style else ""
                prefix = "## " if "Heading 1" in style else ("### " if "Heading" in style else "")
                parts.append(prefix + para.text)
        for table in doc.tables:
            for row in table.rows:
                cells = " | ".join(c.text.strip() for c in row.cells)
                if cells.strip():
                    parts.append(cells)
        return "\n".join(parts) if parts else None
    except ImportError:
        return None
    except Exception:
        return None


def _xlsx_to_csv(path: str) -> str | None:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"# Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) if v is not None else "" for v in row]
                if any(v.strip() for v in vals):
                    # Escape commas and quotes in cell values
                    escaped = []
                    for v in vals:
                        if "," in v or '"' in v or "\n" in v:
                            v = '"' + v.replace('"', '""') + '"'
                        escaped.append(v)
                    parts.append(",".join(escaped))
        return "\n".join(parts) if parts else None
    except ImportError:
        return None
    except Exception:
        return None


def _ipynb_to_text(path: str) -> str | None:
    """Extract code and markdown cells from a Jupyter notebook."""
    try:
        with open(path, encoding="utf-8") as f:
            nb = json.load(f)
        parts = []
        lang = nb.get("metadata", {}).get("kernelspec", {}).get("language", "python")
        parts.append(f"# Jupyter Notebook ({lang})\n")
        for i, cell in enumerate(nb.get("cells", []), 1):
            cell_type = cell.get("cell_type", "code")
            source = "".join(cell.get("source", []))
            if not source.strip():
                continue
            if cell_type == "markdown":
                parts.append(f"## [Markdown Cell {i}]\n{source}")
            elif cell_type == "code":
                parts.append(f"## [Code Cell {i}]\n```{lang}\n{source}\n```")
                # Include cell outputs (text/error only)
                for out in cell.get("outputs", []):
                    out_type = out.get("output_type", "")
                    if out_type in ("stream", "execute_result", "display_data"):
                        text_out = "".join(out.get("text", out.get("data", {}).get("text/plain", [])))
                        if text_out.strip():
                            parts.append(f"_Output:_\n```\n{text_out[:500]}\n```")
                    elif out_type == "error":
                        parts.append(f"_Error:_ {out.get('ename')}: {out.get('evalue', '')}")
        return "\n\n".join(parts) if parts else None
    except Exception:
        return None


def _try_read_as_text(path: str) -> str | None:
    """Try to read an unknown file as UTF-8 text. Returns None for binary files."""
    try:
        with open(path, "r", encoding="utf-8", errors="strict") as f:
            content = f.read(1_000_000)  # max 1MB of text
        # Heuristic: if more than 30% non-printable chars, it's binary
        printable = sum(1 for c in content[:1000] if c.isprintable() or c in "\n\r\t")
        if len(content) > 0 and printable / min(len(content), 1000) < 0.7:
            return None
        return content
    except (UnicodeDecodeError, IOError):
        return None


def _pptx_to_text(path: str) -> str | None:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"\n## Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts) if parts else None
    except ImportError:
        return None
    except Exception:
        return None
